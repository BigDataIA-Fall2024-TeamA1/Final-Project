from fastapi import FastAPI, File, UploadFile, HTTPException
from PyPDF2 import PdfReader
import os
from openai import OpenAI
from openai_module import generate_summary, generate_strategy
from serpapi.google_search import GoogleSearch
from dotenv import load_dotenv
import zipfile
import tempfile
import pandas as pd
import docx
from PyPDF2 import PdfReader
from typing import Union
import easyocr
import logging
from easyocr import Reader
from pinecone import Pinecone, ServerlessSpec
import time
from pydantic import BaseModel

# 加载环境变量
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")

# 初始化 Pinecone 和 OpenAI
pc = Pinecone(api_key=PINECONE_API_KEY)
pinecone_index = pc.Index(PINECONE_INDEX_NAME)
openai_client = OpenAI(api_key=OPENAI_API_KEY)

# 初始化日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = ['.txt', '.pdf', '.docx', '.csv', '.xlsx', '.png', '.jpg', '.jpeg', '.py', '.zip', '.pdb', '.pptx']

app = FastAPI()

class SummaryRequest(BaseModel):
    facts: str
    issues: str
    reasoning: str
    decision: str

def get_embedding(text: str):
    """
    使用 OpenAI 嵌入 API 获取文本的嵌入向量。
    """
    response = openai_client.embeddings.create(
        input=text,
        model="text-embedding-ada-002"
    )
    return response.data[0].embedding

def extract_text_from_file(file_path: str) -> Union[str, None]:
    """
    从不同类型的文件中提取文本内容。
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as f:
                reader_pdf = PdfReader(f)
                for page in reader_pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text

        elif ext == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_csv(index=False)

        elif ext == '.xlsx':
            df = pd.read_excel(file_path, engine='openpyxl')
            return df.to_csv(index=False)

        elif ext in ['.png', '.jpg', '.jpeg']:
            reader = Reader(['en'])
            result = reader.readtext(file_path, detail=0)
            return ' '.join(result)

        elif ext == '.py':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.zip':
            return extract_text_from_zip(file_path)

        elif ext == '.pdb':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        else:
            return f"Unsupported file type: {ext}"

    except Exception as e:
        logger.error(f"Error processing file {file_path}: {e}")
        return f"Error processing file: {e}"

def extract_text_from_zip(file_path: str) -> str:
    """
    从 ZIP 文件中提取内容。
    """
    text = ""
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        with tempfile.TemporaryDirectory() as tmpdir:
            zip_ref.extractall(tmpdir)
            for root, _, files in os.walk(tmpdir):
                for file in files:
                    file_path_inner = os.path.join(root, file)
                    _, ext_inner = os.path.splitext(file_path_inner)
                    ext_inner = ext_inner.lower()
                    if ext_inner in SUPPORTED_EXTENSIONS:
                        extracted_text = extract_text_from_file(file_path_inner)
                        text += f"Extracted from {file}:\n{extracted_text}\n"
                    else:
                        text += f"Skipped unsupported file: {file}\n"
    return text

@app.post("/upload_file_summary")
async def upload_file_summary(file: UploadFile = File(...)):
    """
    通过上传文件生成摘要。
    """
    try:
        # 保存临时文件
        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            tmp_file.write(file.file.read())
            tmp_file_path = tmp_file.name

        # 从文件中提取文本
        file_content = extract_text_from_file(tmp_file_path)
        if not file_content:
            raise HTTPException(status_code=400, detail="Failed to extract text from file.")
        logger.info("Extracted text from uploaded file.")

        # 获取文件内容的嵌入
        user_embedding = get_embedding(file_content)
        logger.info("Generated embedding for uploaded file content.")

        # 查询 Pinecone 以获取最相似的案例
        search_results = pinecone_index.query(vector=user_embedding, top_k=1, include_metadata=True)
        logger.info(f"Search results: {search_results}")

        if "matches" not in search_results or len(search_results["matches"]) == 0:
            raise HTTPException(status_code=404, detail="No similar cases found.")

        # 获取最佳匹配案例
        best_match = search_results["matches"][0]
        metadata = best_match.get("metadata", {})
        case_plain_text = metadata.get("plain_text", "").strip()

        if not case_plain_text:
            raise HTTPException(status_code=400, detail="Similar case content is empty.")

        # 构建 OpenAI 请求
        prompt = f"""
        You are tasked with generating a summary for the following legal case.
        Context from a similar case:
        {case_plain_text}

        Uploaded file content:
        {file_content}

        Generate a detailed and structured summary combining both contexts.
        """
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful AI for legal analysis."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1500,
            temperature=0.5
        )
        summary = response.choices[0].message.content
        return {"summary": summary}

    except Exception as e:
        logger.exception("Error generating summary from file.")
        raise HTTPException(status_code=500, detail=f"Error generating summary from file: {str(e)}")